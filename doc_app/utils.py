

import fitz
import docx
import pptx
from predibase import Predibase

from doc_app.dependencies import get_settings


def extract_text_from_file(file_path, content_type):
    if content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = docx.Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        ppt = pptx.Presentation(file_path)
        text = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif content_type == "application/pdf":
        doc = fitz.open(file_path)
        text = "\n".join([page.get_text() for page in doc])
    return text


async def get_summary(text):
    settings = get_settings()
    pb = Predibase(api_token=settings.predibase_api_key)
    input_prompt = f"<s>[INST] The following passage is content from a news report. Please summarize this passage in three sentence or. \n  {text}. \n Summary: [/INST] "

    lorax_client = pb.deployments.client("llama-3-1-8b-instruct")
    summary_text = lorax_client.generate(input_prompt, max_new_tokens=100).generated_text
    return summary_text
